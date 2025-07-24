#!/usr/bin/env python

import argparse
import asyncio
from asyncio import Task, TaskGroup
import edge_tts
from edge_tts.typing import VoicesManagerFind
import pymupdf
from functools import partial
from itertools import count
import locale
from loguru import logger
from operator import itemgetter
from pathlib import Path
from pptx import Presentation
from pptx.slide import Slide
from string import Template
import subprocess
import tempfile
from typing import Any, TypeGuard

def get_note_from_slide(slide: Slide) -> str | None:
    if not slide.has_notes_slide:
        return None
    
    if not slide.notes_slide.notes_text_frame:
        return None

    notes_text: str = slide.notes_slide.notes_text_frame.text
    if len(notes_text) == 0:
        return None

    return notes_text

def get_notes_from_ppt_file(ppt_file_path: Path) -> list[str | None]:
    prs = Presentation(str(ppt_file_path))
    notes = list(map(get_note_from_slide, prs.slides))
    return notes

async def convert_page_to_image(page: pymupdf.Page,
                          output_file_path: Path,
                          dpi: int) -> Path:
    loop = asyncio.get_running_loop()
    pix: pymupdf.Pixmap = await loop.run_in_executor(None, partial(page.get_pixmap, dpi=dpi))
    loop.run_in_executor(None, pix.save, output_file_path)
    logger.info('Generate Image file from PDF in `{output_file_path}`', output_file_path=output_file_path)

    return output_file_path

async def convert_ppt_to_image(ppt_file_path: Path,
                               output_dir: Path,
                               dpi: int,
                               output_filename: Template,
                               soffice_file_path: Path,
                               encoding: str,
                               pages: list[int] | None = None) -> list[Path]:

    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir_path = Path(tmp_dir)
        output = subprocess.run([soffice_file_path,
                                 '--headless',
                                 '--invisible',
                                 '--convert-to', 'pdf',
                                 '--outdir', tmp_dir_path,
                                 ppt_file_path],
                                check=True, capture_output=True, encoding=encoding)

        if len(output.stderr) > 0:
            raise subprocess.CalledProcessError(output.returncode, output.args, output.stdout, output.stderr)

        pdf_file_path = tmp_dir_path / f'{ppt_file_path.stem}.pdf'
        logger.info('Generate PDF from PPTX in `{pdf_file_path}`', pdf_file_path=pdf_file_path)

        tasks: list[Task[Path]] = list()
        with pymupdf.open(pdf_file_path) as pdf:
            async with TaskGroup() as tg:
                for index, page in enumerate(iter(pdf), start=1):
                    if (pages is not None) and (index not in pages):
                        continue

                    output_file_path = output_dir / output_filename.substitute(index=index)
                    tasks.append(tg.create_task(convert_page_to_image(page, output_file_path, dpi)))

    return list(map(lambda t: t.result(), tasks))

async def convert_note_to_audio(note: str,
                          output_file_path: Path,
                          voice: str) -> Path:
    communicate = edge_tts.Communicate(note, voice)
    await communicate.save(str(output_file_path))
    logger.info('Generate Audio file from note in `{output_file_path}`', output_file_path=output_file_path)

    return output_file_path

async def convert_notes_to_audio(notes: list[str],
                                 output_dir: Path,
                                 output_filename: Template,
                                 voice: str) -> list[Path]:
    tasks: list[Task[Path]] = list()
    async with TaskGroup() as tg:
        for index, note in enumerate(notes, start=1):
            output_file_path = output_dir / output_filename.substitute(index=index)
            tasks.append(tg.create_task(convert_note_to_audio(note, output_file_path, voice)))

    return list(map(lambda task: task.result(), tasks))

def convert_video(image_file_path: Path,
                  audio_file_path: Path,
                  output_file_path: Path,
                  ffmpeg_file_path: Path,
                  encoding: str) -> Path:
    _output = subprocess.run([ffmpeg_file_path,
                             '-loop', '1',
                             '-i', image_file_path,
                             '-i', audio_file_path,
                             '-c:v', 'libx264',
                             '-c:a', 'copy',
                             '-shortest',
                             '-y',
                             output_file_path],
                            check=True, capture_output=True, encoding=encoding)

    logger.info('Generate Video file from Image file and Audio file in `{output_file_path}`', output_file_path=output_file_path)
    return output_file_path

def convert_videos(image_file_paths: list[Path],
                  audio_file_paths: list[Path],
                  output_dir: Path,
                  output_filename: Template,
                  ffmpeg_file_path: Path,
                  encoding: str) -> list[Path]:
    result: list[Path] = list()

    for index, image_file_path, audio_file_path in zip(count(1), image_file_paths, audio_file_paths):
        output_file_path = output_dir / output_filename.substitute(index=index)
        p = convert_video(image_file_path, audio_file_path, output_file_path, ffmpeg_file_path, encoding)
        result.append(p)

    return result

def concat_videos(video_file_paths: list[Path],
                  output_file_path: Path,
                  ffmpeg_file_path: Path,
                  encoding: str) -> Path:
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir_path = Path(tmp_dir)
        concat_file_path = tmp_dir_path / 'concat.txt'
        with concat_file_path.open(mode='w') as f:
            for p in video_file_paths:
                f.write(f"file '{p.resolve()}'\n")
        _output = subprocess.run([ffmpeg_file_path,
                                 '-f', 'concat',
                                 '-safe', '0',
                                 '-i', concat_file_path,
                                 '-c:v', 'copy',
                                 '-c:a', 'aac',
                                 '-ar', '48000',
                                 '-y',
                                 output_file_path],
                                check=True, capture_output=True, encoding=encoding)

        logger.info('Concat Video file from several Video files in `{output_file_path}`', output_file_path=output_file_path)
        return output_file_path

def has_note(t: tuple[int, str | None]) -> TypeGuard[tuple[int, str]]:
    _index, note = t
    return (note is not None) and (len(note) > 0)

async def main_process(ppt_file_path: Path,
                       output_file_path: Path,
                       soffice_file_path: Path,
                       ffmpeg_file_path: Path,
                       dpi: int,
                       voice: str,
                       encoding: str) -> Path:
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_dir_path = Path(tmp_dir)
        notes = get_notes_from_ppt_file(ppt_file_path)

        available_pages_and_notes: list[tuple[int, str]] = list(filter(has_note, enumerate(notes, start=1)))
        available_pages: list[int] = list(map(itemgetter(0), available_pages_and_notes))
        available_notes: list[str] = list(map(itemgetter(1), available_pages_and_notes))

        audio_dir_path = tmp_dir_path / 'audios'
        audio_dir_path.mkdir()
        audio_file_paths = await convert_notes_to_audio(notes=available_notes,
                                                        output_dir=audio_dir_path,
                                                        output_filename=Template('note-${index}.aac'),
                                                        voice=voice)

        image_dir_path = tmp_dir_path / 'images'
        image_dir_path.mkdir()
        image_file_paths = await convert_ppt_to_image(ppt_file_path=ppt_file_path,
                                                      output_dir=image_dir_path,
                                                      pages=available_pages,
                                                      dpi=dpi,
                                                      output_filename=Template('page-${index}.png'),
                                                      soffice_file_path=soffice_file_path,
                                                      encoding=encoding)

        video_dir_path = tmp_dir_path / 'videos'
        video_dir_path.mkdir()
        video_file_paths = convert_videos(image_file_paths=image_file_paths,
                                          audio_file_paths=audio_file_paths,
                                          output_dir=video_dir_path,
                                          output_filename=Template('video-${index}.mp4'),
                                          ffmpeg_file_path=ffmpeg_file_path,
                                          encoding=encoding)

        result = concat_videos(video_file_paths=video_file_paths,
                               output_file_path=output_file_path,
                               ffmpeg_file_path=ffmpeg_file_path,
                               encoding=encoding)

        return result

async def convert(args: argparse.Namespace) -> Path:
    result = await main_process(ppt_file_path=args.infile,
                        output_file_path=args.outfile,
                        soffice_file_path=args.soffice_file_path,
                        ffmpeg_file_path=args.ffmpeg_file_path,
                        dpi=args.dpi,
                        voice=args.voice,
                        encoding=args.encoding)
    return result

def pretty_format(obj: dict[Any, Any] | list[Any] | str, depth: int = 0) -> str:
    result: list[str] = list()

    if isinstance(obj, dict):
        for k, v in obj.items():
            k_str = pretty_format(k, depth)
            if isinstance(v, str):
                result.append(': '.join([k_str, str(v)]))
            else:
                result.append(f'{k_str}:')
                v_str = pretty_format(v, depth+1)
                result.append(v_str)
        if depth == 0:
            result.append('')
    elif isinstance(obj, list):
        for item in obj:
            result.append(pretty_format(item, depth))
    else:
        indent = ' ' * 2 * depth
        result.append(f'{indent}{str(obj)}')

    return '\n'.join(result)

async def list_voices(args: argparse.Namespace) -> None:
    params: VoicesManagerFind = {}

    if args.language != 'all':
        params['Language'] = args.language

    if args.locale != 'all':
        params['Locale'] = args.locale

    if args.gender != 'all':
        params['Gender'] = args.gender.capitalize()

    voices_manager = await edge_tts.VoicesManager.create()
    voices = voices_manager.find(**params)
    voices.sort(key=itemgetter('ShortName'))

    if not args.detail:
        for name in map(itemgetter('ShortName'), voices):
            print(name)
        return

    print(pretty_format(voices))

def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Convert PPT file(s) to one Video file")
    subparsers = parser.add_subparsers(required=True)

    parser_convert = subparsers.add_parser('convert')
    parser_convert.add_argument('-i', '--infile', type=Path, help='add PPT file(s)', required=True)
    parser_convert.add_argument('outfile', type=Path, help='set output video filename')
    parser_convert.add_argument('--soffice-file-path', type=Path, default=Path('/usr/bin/soffice'))
    parser_convert.add_argument('--ffmpeg-file-path', type=Path, default=Path('/usr/bin/ffmpeg'))
    parser_convert.add_argument('--dpi', type=int, default=75)
    parser_convert.add_argument('--voice', type=str, default='zh-CN-XiaoxiaoNeural')
    parser_convert.add_argument('--encoding', type=str, default=locale.getpreferredencoding())
    parser_convert.set_defaults(func=convert)

    parser_list_voices = subparsers.add_parser('list-voices')
    parser_list_voices.add_argument('--language', type=str, default='all')
    parser_list_voices.add_argument('--locale', type=str, default='all')
    parser_list_voices.add_argument('--gender', type=str, default='all')
    parser_list_voices.add_argument('--detail', action='store_true', default=False)
    parser_list_voices.set_defaults(func=list_voices)

    return parser.parse_args()

async def main() -> None:
    args = parse_args()

    await args.func(args)

if __name__ == '__main__':
    asyncio.run(main())
